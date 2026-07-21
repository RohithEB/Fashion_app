#!/usr/bin/env python3
"""Generate the Maison Ebani / Fashion Showroom presentation deck.

Usage:
    pip install python-pptx
    python scripts/make_ppt.py            # writes docs/Fashion_Showroom_Deck.pptx

Edit the SLIDES data at the bottom of this file and re-run to update the deck.
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------- design tokens

INK        = RGBColor(0x0E, 0x0E, 0x10)   # background
INK_SOFT   = RGBColor(0x1A, 0x1A, 0x1E)   # panels
IVORY      = RGBColor(0xF5, 0xF2, 0xEC)   # primary text
MUTED      = RGBColor(0x9A, 0x97, 0x90)   # secondary text
CHAMPAGNE  = RGBColor(0xC6, 0xA6, 0x64)   # accent
LINE       = RGBColor(0x33, 0x33, 0x38)   # rules / borders

SERIF = "Georgia"          # stands in for Cormorant Garamond
SANS  = "Segoe UI"         # stands in for Inter
MONO  = "Consolas"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.9)


# ------------------------------------------------------------------- primitives

def _blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = INK
    return s


def _box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def _para(tf, text, *, size, color, font=SANS, bold=False, italic=False,
          space_before=0, space_after=6, align=PP_ALIGN.LEFT, line=None, first=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line:
        p.line_spacing = line
    r = p.add_run()
    r.text = text
    f = r.font
    f.size, f.bold, f.italic, f.name = Pt(size), bold, italic, font
    f.color.rgb = color
    return p


def _rule(slide, x, y, w, color=CHAMPAGNE, thick=Pt(1.5)):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _panel(slide, x, y, w, h, fill=INK_SOFT):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.04
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = LINE
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    return sh


def _footer(slide, n, label):
    tf = _box(slide, MARGIN, H - Inches(0.62), W - 2 * MARGIN, Inches(0.3))
    _para(tf, f"Maison Ebani  ·  {label}", size=9, color=MUTED, first=True, space_after=0)
    tf2 = _box(slide, W - MARGIN - Inches(1.0), H - Inches(0.62), Inches(1.0), Inches(0.3))
    _para(tf2, str(n), size=9, color=MUTED, align=PP_ALIGN.RIGHT, first=True, space_after=0)


def _heading(slide, title, kicker=None):
    y = Inches(0.62)
    if kicker:
        tf = _box(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.3))
        _para(tf, kicker.upper(), size=10.5, color=CHAMPAGNE, bold=True, first=True, space_after=0)
        y += Inches(0.36)
    tf = _box(slide, MARGIN, y, W - 2 * MARGIN, Inches(0.7))
    _para(tf, title, size=31, color=IVORY, font=SERIF, first=True, space_after=0)
    _rule(slide, MARGIN, y + Inches(0.66), Inches(1.4))
    return y + Inches(1.0)


# ------------------------------------------------------------------ slide kinds

def slide_title(prs, s, n):
    sl = _blank(prs)
    _rule(sl, MARGIN, Inches(2.35), Inches(2.2))
    tf = _box(sl, MARGIN, Inches(2.6), W - 2 * MARGIN, Inches(1.6))
    _para(tf, s["title"], size=54, color=IVORY, font=SERIF, first=True, space_after=10)
    _para(tf, s["subtitle"], size=19, color=CHAMPAGNE, font=SERIF, italic=True, space_after=16)
    _para(tf, s["blurb"], size=13, color=MUTED, line=1.5)
    tf = _box(sl, MARGIN, H - Inches(1.25), W - 2 * MARGIN, Inches(0.5))
    _para(tf, s["meta"], size=10.5, color=MUTED, first=True, space_after=0)
    return sl


def slide_section(prs, s, n):
    sl = _blank(prs)
    tf = _box(sl, MARGIN, Inches(2.9), W - 2 * MARGIN, Inches(0.5))
    _para(tf, s["kicker"].upper(), size=11, color=CHAMPAGNE, bold=True, first=True, space_after=8)
    tf = _box(sl, MARGIN, Inches(3.35), W - 2 * MARGIN, Inches(1.2))
    _para(tf, s["title"], size=42, color=IVORY, font=SERIF, first=True, space_after=10)
    if s.get("blurb"):
        _para(tf, s["blurb"], size=13.5, color=MUTED, line=1.45)
    _rule(sl, MARGIN, Inches(3.24), Inches(1.4))
    _footer(sl, n, s["kicker"])
    return sl


def slide_bullets(prs, s, n):
    sl = _blank(prs)
    y = _heading(sl, s["title"], s.get("kicker"))
    tf = _box(sl, MARGIN, y, W - 2 * MARGIN, H - y - Inches(0.9))
    first = True
    for b in s["bullets"]:
        if isinstance(b, tuple):
            lead, rest = b
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.space_before, p.space_after, p.line_spacing = Pt(0), Pt(11), 1.3
            r = p.add_run(); r.text = lead + "  "
            r.font.size, r.font.bold, r.font.name = Pt(15), True, SANS
            r.font.color.rgb = CHAMPAGNE
            r2 = p.add_run(); r2.text = rest
            r2.font.size, r2.font.name = Pt(15), SANS
            r2.font.color.rgb = IVORY
        else:
            _para(tf, b, size=15, color=IVORY, space_after=11, line=1.3, first=first)
        first = False
    _footer(sl, n, s.get("kicker", "Fashion Showroom"))
    return sl


def slide_mono(prs, s, n):
    """Heading + monospaced ASCII diagram in a panel."""
    sl = _blank(prs)
    y = _heading(sl, s["title"], s.get("kicker"))
    body = s["mono"]
    lines = body.count("\n") + 1
    size = s.get("size", 11)
    h = min(H - y - Inches(0.95), Inches(0.235 * lines * (size / 11) + 0.5))
    _panel(sl, MARGIN, y, W - 2 * MARGIN, h)
    tf = _box(sl, MARGIN + Inches(0.3), y + Inches(0.24), W - 2 * MARGIN - Inches(0.6), h - Inches(0.4))
    first = True
    for ln in body.split("\n"):
        _para(tf, ln or " ", size=size, color=IVORY, font=MONO, space_after=0, line=1.15, first=first)
        first = False
    if s.get("note"):
        tf = _box(sl, MARGIN, y + h + Inches(0.2), W - 2 * MARGIN, Inches(0.6))
        _para(tf, s["note"], size=12.5, color=MUTED, first=True, line=1.35)
    _footer(sl, n, s.get("kicker", "Fashion Showroom"))
    return sl


def slide_columns(prs, s, n):
    sl = _blank(prs)
    y = _heading(sl, s["title"], s.get("kicker"))
    cols = s["columns"]
    gap = Inches(0.35)
    cw = (W - 2 * MARGIN - gap * (len(cols) - 1)) / len(cols)
    ch = H - y - Inches(1.0)
    for i, c in enumerate(cols):
        x = MARGIN + i * (cw + gap)
        _panel(sl, x, y, cw, ch)
        tf = _box(sl, x + Inches(0.28), y + Inches(0.26), cw - Inches(0.56), ch - Inches(0.5))
        _para(tf, c["head"], size=17, color=CHAMPAGNE, font=SERIF, first=True, space_after=4)
        if c.get("sub"):
            _para(tf, c["sub"], size=10.5, color=MUTED, space_after=12)
        for it in c["items"]:
            _para(tf, "·  " + it, size=12.5, color=IVORY, space_after=7, line=1.25)
    _footer(sl, n, s.get("kicker", "Fashion Showroom"))
    return sl


def slide_table(prs, s, n):
    sl = _blank(prs)
    y = _heading(sl, s["title"], s.get("kicker"))
    rows, cols = len(s["rows"]) + 1, len(s["head"])
    th = min(H - y - Inches(0.95), Inches(0.42) * rows + Inches(0.1))
    shp = sl.shapes.add_table(rows, cols, MARGIN, y, W - 2 * MARGIN, th)
    tbl = shp.table
    tbl.first_row = True
    widths = s.get("widths")
    if widths:
        total = W - 2 * MARGIN
        for i, frac in enumerate(widths):
            tbl.columns[i].width = Emu(int(total * frac))
    for c, txt in enumerate(s["head"]):
        cell = tbl.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = INK
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.12)
        p = cell.text_frame.paragraphs[0]
        r = p.add_run(); r.text = txt
        r.font.size, r.font.bold, r.font.name = Pt(11.5), True, SANS
        r.font.color.rgb = CHAMPAGNE
    for ri, row in enumerate(s["rows"], start=1):
        for ci, txt in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK_SOFT if ri % 2 else INK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.12)
            p = cell.text_frame.paragraphs[0]
            p.line_spacing = 1.1
            r = p.add_run(); r.text = txt
            r.font.size, r.font.name = Pt(s.get("size", 11)), SANS
            r.font.color.rgb = IVORY if ci == 0 else MUTED
    _footer(sl, n, s.get("kicker", "Fashion Showroom"))
    return sl


KINDS = {
    "title": slide_title, "section": slide_section, "bullets": slide_bullets,
    "mono": slide_mono, "columns": slide_columns, "table": slide_table,
}


# ------------------------------------------------------------------- the content

SLIDES = [
 {"kind": "title",
  "title": "Maison Ébani",
  "subtitle": "Fashion Showroom — a guided luxury selling experience",
  "blurb": "A salesperson's phone privately drives a large in-store display.\n"
           "Not screen mirroring — event-based Live Presentation Synchronization.",
  "meta": "Proof of Concept  ·  Flutter mobile + display  ·  Node/SQLite backend  ·  Next.js CMS"},

 {"kind": "bullets", "kicker": "The problem", "title": "“Here, look at my phone.”",
  "bullets": [
    ("The moment breaks the sale.", "A five-inch screen, held sideways, with the associate's private notes, prices-in-progress and UI chrome all visible to the guest."),
    ("Mirroring is worse.", "Streaming the phone shows every search, every mis-tap, every hesitation — and needs bandwidth a boutique WiFi rarely has."),
    ("Luxury retail needs two views.", "A private, information-rich cockpit for the associate; a large, editorial, distraction-free stage for the customer."),
    ("And it must work offline.", "In-store WiFi drops. The experience cannot."),
  ]},

 {"kind": "bullets", "kicker": "The solution", "title": "One product, two screens, zero leakage",
  "bullets": [
    ("Private by default.", "The associate browses, searches and filters freely. None of it reaches the display."),
    ("Presentation on demand.", "“Show on Screen” starts the synchronized presentation — and only then."),
    ("Events, not video.", "Tiny WebSocket messages (show_product, change_colour, zoom, scroll) describe the interaction; the display renders it from its own cached catalogue."),
    ("Offline-first.", "Both apps ship a 37-product catalogue snapshot and draw ~97% of imagery on-device. Phone and display just need the same WiFi."),
    ("Production-shaped.", "Layered backend, repository seams, frozen protocol, e2e tests — a POC that scales into a product."),
  ]},

 {"kind": "mono", "kicker": "Architecture", "title": "How the pieces fit", "size": 11.5,
  "mono": " MOBILE (salesperson)                          DISPLAY (customer screen)\n"
          " - browse privately (no sync)                  - phase state machine (FSM)\n"
          " - \"Show on Screen\" ----+                      - renders from its own catalogue\n"
          " - live controls -------|   WebSocket events   - applies each event:\n"
          "   colour / size /      |   show_product,        show_product   -> presenting\n"
          "   image / zoom /       |   change_*, zoom,      change_colour  -> new variant\n"
          "   scroll / video       |   scroll, clear,       zoom / scroll  -> transform\n"
          "                        +-> payment_success -->  payment_success-> Thank-You\n"
          "\n"
          "        +---------------- one of two transports -------------------+\n"
          "        |  BOX MODE (default): display hosts the LAN WS server     |\n"
          "        |  BACKEND MODE:       Node server relays controller->display|\n"
          "        +----------------------------------------------------------+\n"
          "\n"
          " CMS (Next.js) --reads/writes--> SQLite <--reads/writes-- Node server",
  "note": "Backend is strictly layered: repositories (SQL) → services (business logic) → routes + ws. "
          "Frontend is feature-first: each feature owns its controller and screens."},

 {"kind": "columns", "kicker": "Applications", "title": "Four independently runnable apps",
  "columns": [
    {"head": "mobile_app", "sub": "Flutter · portrait · associate's phone",
     "items": ["Login & salesperson profile", "QR pairing to a display", "Search, filter, colour explorer",
               "Product detail + live controls", "Saved outfits, checkout", "Private AI talking point"]},
    {"head": "display_app", "sub": "Flutter · dark, full-bleed · box/TV",
     "items": ["Phase state machine", "Pairing QR + welcome", "Editorial product presentation",
               "Gallery, video, colourways", "Thank-you + idle return", "Hosts the LAN server in box mode"]},
    {"head": "server", "sub": "Node 20 · Express · ws · SQLite · :3000",
     "items": ["Auth + bearer tokens", "Catalogue, search, filters", "Recommendations & talking points",
               "Cart, checkout, orders", "WS relay + session rules", "Playwright e2e suite"]},
    {"head": "cms", "sub": "Next.js 15 · React 19 · MUI 6 · :4000",
     "items": ["Product CRUD with images", "AI enrichment via Claude", "Cloudflare R2 uploads",
               "Salespeople management", "Dashboard analytics", "Reads the same SQLite file"]},
  ]},

 {"kind": "mono", "kicker": "Experience", "title": "The session, end to end", "size": 12.5,
  "mono": "Login\n"
          "  -> [optional] Salesperson profile setup\n"
          "  -> Pair display (scan the QR on screen)\n"
          "  -> Customer details (optional, skippable)\n"
          "  -> Top-6 recommendations\n"
          "  -> Browse / search / filter        (private — nothing on screen)\n"
          "  -> Product detail -> \"Show on Screen\"   <-- presentation begins\n"
          "  -> Live controls: colour · size · image · zoom · scroll · gallery · video\n"
          "  -> Saved outfits -> Checkout -> Payment success\n"
          "  -> End session -> display returns to a fresh pairing QR",
  "note": "Display phases: splash → advertisement → waiting (QR) → connecting → loading → welcome → "
          "presenting | catalogue | cart | checkout → thankYou → (countdown) → waiting."},

 {"kind": "bullets", "kicker": "The differentiator", "title": "Live Presentation Synchronization",
  "bullets": [
    ("The phone screen is never streamed.", "A pinch-to-zoom emits zoom { scale, focalX, focalY }; the display applies the same transform to its own cached image."),
    ("Synchronized set.", "Show / hide product · colour variant · size · image swipe · pinch-zoom · pan · detail-scroll · gallery grid · video play-pause-seek · reset view."),
    ("A pure reducer on both sides.", "ProductPresentation.applyEvent(WsEvent) runs identically on mobile and display — the mobile mirrors exactly what the guest sees."),
    ("Cheap on the wire.", "Payloads are ids plus transform parameters; high-frequency gestures are throttled and coalesced before emission."),
    ("Always-visible truth.", "The associate always has a “now showing” indicator — no guessing what the customer can see."),
  ]},

 {"kind": "table", "kicker": "Protocol", "title": "The frozen WebSocket contract",
  "widths": [0.2, 0.44, 0.36], "size": 11,
  "head": ["Direction", "Messages", "Rule"],
  "rows": [
    ["Controller → server", "pair, activity, keep_alive", "Binds one controller to one display; resets the idle timer"],
    ["Relayed commands", "show_catalog, show_cart, show_product, show_details, show_related, show_media, zoom, clear", "Allow-listed by RELAY_TYPES — anything else is rejected"],
    ["Server → clients", "display_registered, paired, session_warning, session_end, error", "Server is authoritative; the display never trusts the mobile directly"],
    ["Session lifecycle", "10-min idle → session_warning → grace → session_end", "Ending frees the display and re-issues a fresh QR"],
    ["Box mode", "Richer WsEvent vocabulary, no relay translation", "The display hosts the socket directly"],
  ]},

 {"kind": "table", "kicker": "Backend", "title": "HTTP API surface",
  "widths": [0.34, 0.66], "size": 11,
  "head": ["Endpoint", "Purpose"],
  "rows": [
    ["/auth/register · /login · /logout · /me", "Salesperson auth, bearer token"],
    ["/products · /products/:id · /:id/similar", "List, search, filter, sort, page; full detail; cross-product “show similar”"],
    ["/categories · /filters", "Facets — categories, colours, sizes, price range"],
    ["/recommendations · /talking-point", "Profile-matched picks; private coaching cue for the associate"],
    ["/customers · /customers/options", "Capture the session guest; onboarding choice lists"],
    ["/cart/:sessionId  (+ /default, /checkout)", "Server-side cart, auto-displayed item, order placement"],
    ["/orders/:id · /journey · /health", "Order detail, journey events, health probe"],
    ["/media/ph?w&h&bg&fg&text", "Generated SVG placeholder imagery — the offline image strategy"],
  ]},

 {"kind": "columns", "kicker": "Feature set", "title": "What is built",
  "columns": [
    {"head": "Presentation & sync",
     "items": ["Show product / clear", "Colour variant, size", "Image swipe + gallery grid",
               "Fullscreen pinch-zoom + pan", "Detail-scroll mirroring", "Product video", "Reset view"]},
    {"head": "Catalogue",
     "items": ["Keyword search", "Category chips", "Colour explorer", "Collapsing header",
               "Pull-to-refresh", "Variants, media, enrichment"]},
    {"head": "Personalisation",
     "items": ["Session-scoped guest profile", "Top-6 scored recommendations",
               "Gender · personality · style hints", "Private on-phone AI talking point",
               "Never shown to the customer"]},
    {"head": "Commerce & admin",
     "items": ["Saved outfits (shortlist)", "Doubles as on-screen selector",
               "Checkout with pre-filled guest", "Order confirmation", "Journey logging",
               "CMS: CRUD, AI, analytics"]},
  ]},

 {"kind": "table", "kicker": "Engineering", "title": "Technology choices",
  "widths": [0.22, 0.46, 0.32], "size": 11,
  "head": ["Layer", "Stack", "Why"],
  "rows": [
    ["Both apps", "Flutter 3.41 / Dart 3.11, Material 3", "One codebase for phone and Android box"],
    ["State & nav", "provider (ChangeNotifier) + go_router guards", "One controller per feature, DI via MultiProvider"],
    ["Look", "Cormorant Garamond + Inter, Material Symbols", "Editorial luxury, fully tokenised — no hardcoded values"],
    ["Realtime", "web_socket_channel; ws on the server", "Events, not video"],
    ["Backend", "Node 20 · Express · better-sqlite3 · qrcode", "Single process, single file of truth"],
    ["CMS", "Next.js 15 · React 19 · MUI 6 · Anthropic SDK · R2", "Back-office editing and AI enrichment"],
    ["Quality", "Playwright e2e (API + WS), db:check smoke", "The contract is tested, not assumed"],
  ]},

 {"kind": "columns", "kicker": "Deployment", "title": "Two runtime modes, one binary",
  "columns": [
    {"head": "Box mode  ·  default", "sub": "BACKEND=false — no internet required",
     "items": ["The display hosts HTTP + WS on :8080",
               "Its QR carries its own LAN IP + pairing token",
               "Both apps read a bundled 37-product snapshot",
               "Placeholder imagery is drawn on-device (~97% of images)",
               "Recommendations and the talking point computed locally",
               "Auth / checkout / customers are mocked — no persistence"]},
    {"head": "Backend mode", "sub": "BACKEND=true, BACKEND_HOST=<lan-ip>",
     "items": ["Catalogue, auth, customers, cart, checkout over HTTP",
               "Realtime relayed through the server",
               "Everything persists to SQLite and appears in the CMS",
               "Requires the Node server reachable on the LAN",
               "In-app Server settings sheet overrides the compile-time default at runtime"]},
  ]},

 {"kind": "mono", "kicker": "Run it", "title": "From clone to live demo", "size": 12,
  "mono": "# 1) Dependencies\n"
          "cd server && npm install\n"
          "cd ../cms && npm install\n"
          "cd ../frontend/mobile_app && flutter pub get\n"
          "cd ../display_app        && flutter pub get\n"
          "\n"
          "# 2) Offline demo (default) — no backend needed\n"
          "cd frontend/display_app && flutter run    # start the display FIRST\n"
          "cd frontend/mobile_app  && flutter run    # scan the QR on the display\n"
          "\n"
          "# 3) Persisted demo\n"
          "cd server && npm start                    # :3000, seeds on first run\n"
          "cd cms    && npm run dev                  # :4000\n"
          "flutter run --dart-define=BACKEND=true --dart-define=BACKEND_HOST=192.168.1.5",
  "note": "npm test runs the Playwright API + WS suite. Both devices must be on the same WiFi for a live session."},

 {"kind": "bullets", "kicker": "Honesty", "title": "Known limitations",
  "bullets": [
    ("Box mode does not persist.", "Auth, checkout and guest capture are mocked — orders never reach the CMS. Use backend mode when data must survive."),
    ("The Node server does not run on the Android box yet.", "Node + better-sqlite3 need Termux (fragile) or a Dart re-implementation; today backend mode needs a laptop on the LAN."),
    ("Some guest fields are session-scoped.", "DOB, occupation, budget, brands and notes have no backend column yet — modelled for a future API."),
    ("~17 catalogue images are cloud URLs.", "The other ~593 are on-device placeholders and work fully offline."),
    ("The two Flutter apps duplicate shared files.", "Models, realtime and widgets are kept in sync by hand — extraction into a package is planned."),
    ("ingest/ is scaffolded but unwired.", "It needs INGEST_URL and a source feed."),
  ]},

 {"kind": "columns", "kicker": "Roadmap", "title": "Where it goes next",
  "columns": [
    {"head": "Near term",
     "items": ["On-device SQLite so box mode persists and syncs when online",
               "Backend columns + PATCH /customers/:id for the full guest profile",
               "Salesperson profile API — email, avatar, role",
               "Bundle the remaining images for a 100% offline catalogue"]},
    {"head": "Medium term",
     "items": ["Run the server on the box itself (Termux node:sqlite, or a Dart HTTP/WS server inside the display app)",
               "Wire ingest/ for scheduled catalogue refresh",
               "Extract the duplicated Flutter foundation into a shared package",
               "Real AI talking points with an offline template fallback"]},
    {"head": "Longer term",
     "items": ["Journey replay and conversion reporting per associate",
               "Multi-display, multi-session per store",
               "Role-based access control and audit trails",
               "Automated widget + integration coverage for both apps"]},
  ]},

 {"kind": "section", "kicker": "In one line",
  "title": "The guest sees the piece. Never the phone.",
  "blurb": "Private cockpit · editorial stage · event sync · offline by default."},
]


def build(out_path: Path) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for i, s in enumerate(SLIDES, start=1):
        KINDS[s["kind"]](prs, s, i)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    root = Path(__file__).resolve().parent.parent
    p = build(root / "docs" / "Fashion_Showroom_Deck.pptx")
    print(f"wrote {p}  ({len(SLIDES)} slides)")
